
import csv
import gc
import imageio_ffmpeg
import io
import json
import logging
import os
import soundfile as sf
import subprocess
import tarfile
import tempfile
import traceback
import whisper
import xml.etree.ElementTree as ET
import zipfile
# ======================================
from base64 import b64encode
from datetime import datetime
from docx import Document
from moviepy.editor import AudioFileClip
from moviepy.editor import VideoFileClip
from openpyxl import load_workbook
from pptx import Presentation
from PyPDF2 import PdfReader
from yt_dlp import YoutubeDL
# from pytube import YouTube
# ======================================

# import io
# import os
# import csv
# import zipfile
# import tarfile
# import json

# import os
# import gc
# import json
# import tempfile
# import traceback
# import subprocess
# import soundfile as sf
# from moviepy.editor import VideoFileClip
# from yt_dlp import YoutubeDL
# import whisper
# import logging
# from datetime import datetime
# import tempfile
# import gc
# import soundfile as sf
# import subprocess
# import whisper

# import logging
# from datetime import datetime
# import tempfile
# import gc
# import soundfile as sf

# import logging
# from datetime import datetime

# import gc
# import soundfile as sf

# # ============================================
# import xml.etree.ElementTree as ET
# from PyPDF2 import PdfReader
# from docx import Document
# from pptx import Presentation
# from openpyxl import load_workbook
# # ============================================
# import logging
# import tempfile
# import whisper
# import subprocess
# from base64 import b64encode
# # from pytube import YouTube
# from yt_dlp import YoutubeDL
# # ============================================
# import imageio_ffmpeg
# from moviepy.editor import VideoFileClip
# from moviepy.editor import AudioFileClip
# ============================================
import warnings
warnings.filterwarnings("ignore", message="FP16 is not supported on CPU*")
# ============================================#============================================
from models.LLM_GEMINI import *
# from LLM_GEMINI import summarize_frames_with_gemini, analyze_video_frames, classify_video_content_type
# ============================================#============================================
def route_audio_handler(filename, decoded_bytes, query):
    query_lower = query.lower()

    if "chord" in query_lower or "guitar" in query_lower or "lyric" in query_lower:
        return handle_audio_lyrics_chords(filename, decoded_bytes)

    # elif "lyric" in query_lower:
    #     return handle_audio_lyrics_only(filename, decoded_bytes)
    # elif "summary" in query_lower or "explain" in query_lower:
    #     return handle_audio_summary(filename, decoded_bytes)
    else:
        # default to basic audio handler
        return handle_audio(filename, decoded_bytes)

# # =============================================================================
# 📦 If You Want One Recommendation
# For local audio/video file understanding (speech + visuals + text reasoning):
#     ✅ Use Whisper + Ollama (LLaMA3 or Mistral) for audio
#     ✅ Use ffmpeg + BLIP2 or MiniGPT-4 for video frames
#     Wrap everything in Python with a Streamlit or FastAPI interface
# # =============================================================================
def handle_audio_lyrics_chords(filename, decoded_bytes):
    def log(msg):
        print(f"[{datetime.now().strftime('%H:%M:%S')}] {msg}")

    logging.basicConfig(level=logging.INFO)


    audio_path = None
    converted_audio_path = None

    try:
        log(f"📥 Saving audio file: {filename}")

        with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as tmp_audio:
            tmp_audio.write(decoded_bytes)
            tmp_audio.flush()
            audio_path = tmp_audio.name

        converted_audio_path = audio_path.replace(".mp3", "_16k.wav")

        log("🏚 Converting audio to 16kHz mono WAV using ffmpeg...")
        # subprocess.run([
        #     "ffmpeg", "-y", "-t", "120", "-i", audio_path,
        #     "-ac", "1", "-ar", "16000", "-sample_fmt", "s16",
        #     "-f", "wav", converted_audio_path
        # ], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)

        result = subprocess.run([
            "ffmpeg", "-y", "-i", audio_path,
            "-ac", "1", "-ar", "16000", "-sample_fmt", "s16",
            "-f", "wav", converted_audio_path
        ], capture_output=True, text=True)

        if result.returncode != 0:
            print("❌ FFmpeg failed:", result.stderr)
        else:
            print("✅ FFmpeg success:", result.stdout)


        log("🧠 Loading Whisper model (medium/base)...")
        model = whisper.load_model("base")

        log("🧺 Validating audio...")
        data, samplerate = sf.read(converted_audio_path)
        if len(data) < samplerate:
            raise Exception("Audio too short for transcription")

        log("🗣 Transcribing...")
        result = model.transcribe(converted_audio_path)
        transcript = result.get("text", "").strip()
        lang_code = result.get("language", "en")

        log(f">> Transcript (preview): {transcript[:100]}")

        lang_hint = detect_language_hint(filename, transcript, LANGUAGE_HINTS)

        if lang_hint and lang_hint != lang_code:
            log(f">> Retrying transcription with lang hint: {lang_code} → {lang_hint}")
            try:
                result_hint = model.transcribe(converted_audio_path, language=lang_hint)
                new_transcript = result_hint.get("text", "").strip()
                if len(new_transcript.split()) > len(transcript.split()):
                    transcript = new_transcript
                    lang_code = lang_hint
                    log(f">> Improved transcript (preview): {transcript[:100]}")
            except Exception as e:
                log(f"⚠️ Lang hint retry failed: {e}")

        if not transcript:
            return f"\n\n🎧 AUDIO: {filename}\n[⚠️ Transcription failed or invalid.]"

        log("📝 Generating structured lyrics and chords...")

        smart_prompt = (
            "This is a transcript of a song. Format it as a song sheet.\n"
            "Add appropriate line breaks and sections (like Verse, Chorus, Bridge).\n"
            "If possible, infer common chord progressions and place chords above lyrics.\n"
            "Only use chords like [C], [G], [Am], [F], etc. where they make musical sense.\n"
            "If unsure, leave that part without chords.\n\n"
            f"Transcript:\n{transcript[:3000]}"
        )

        final_output = summarize_frames_with_gemini("", smart_prompt, transcript)

        # return f"\n\n🎧 AUDIO: {filename}\n{final_output}"
        from fastapi.responses import JSONResponse

        return JSONResponse(content={
            "type": "audio",
            "filename": filename,
            "summary": final_output[:4000] + "..." if len(final_output) > 4000 else final_output
        })
        print("✅ Response returned to frontend.")


    except Exception as e:
        logging.error("Audio processing failed: %s", str(e))
        return f"[Error processing audio file {filename}: {e}]"

    finally:
        gc.collect()
        for path in [audio_path, converted_audio_path]:
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                    logging.info(f">> [cleanup] Deleted temp file: {path}")
                except Exception as cleanup_error:
                    logging.warning(f">> [cleanup warning] Failed to delete {path}: {cleanup_error}")

# ==========================================================================
# ==========================================================================
def handle_audio(filename, decoded_bytes):
    def log(msg):
        print(f"[{datetime.now().strftime('%H:%M:%S')}] {msg}")

    logging.basicConfig(level=logging.INFO)

    audio_path = None
    converted_audio_path = None

    try:
        log(f"📥 Saving audio file: {filename}")
        
        # 👇 TEMPORARY DEBUG LINE to test incoming audio
        print("✅ Received decoded_bytes:", len(decoded_bytes), "bytes")
        print("🔍 First few bytes (base64):", b64encode(decoded_bytes[:20]).decode())

        with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as tmp_audio:
            tmp_audio.write(decoded_bytes)
            tmp_audio.flush()
            audio_path = tmp_audio.name

        converted_audio_path = audio_path.replace(".mp3", "_16k.wav")

        log("🎚 Converting audio to 16kHz mono WAV using ffmpeg...")
        # subprocess.run([
        #     "ffmpeg", "-y", "-i", audio_path,
        #     "-ac", "1", "-ar", "16000", "-sample_fmt", "s16",
        #     "-f", "wav", converted_audio_path
        # ], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)


        result = subprocess.run([
            "ffmpeg", "-y", "-i", audio_path,
            "-ac", "1", "-ar", "16000", "-sample_fmt", "s16",
            "-f", "wav", converted_audio_path
        ], capture_output=True, text=True)

        if result.returncode != 0:
            print("❌ FFmpeg failed:", result.stderr)
        else:
            print("✅ FFmpeg success:", result.stdout)

        log("🧠 Loading Whisper model (medium/base)...")
        model = whisper.load_model("base")

        log("🧪 Validating audio...")
        data, samplerate = sf.read(converted_audio_path)
        if len(data) < samplerate:
            raise Exception("Audio too short for transcription")

        log("🗣 Transcribing (first pass)...")
        result = model.transcribe(converted_audio_path)
        transcript = result.get("text", "").strip()
        lang_code = result.get("language", "en")
        log(f">> Transcript preview: {transcript[:100]}")

        lang_hint = detect_language_hint(filename, transcript, LANGUAGE_HINTS)
        if lang_hint and lang_hint != lang_code:
            log(f"🗣 Retrying with language hint override: {lang_code} → {lang_hint}")
            try:
                result = model.transcribe(converted_audio_path, language=lang_hint)
                transcript = result.get("text", "").strip()
                lang_code = lang_hint
                log(f">> Overridden transcript preview: {transcript[:100]}")
            except Exception as e:
                log(f"⚠️ Language override failed: {e}")

        if is_transcript_repetitive(transcript):
            log("⚠️ Transcript is repetitive. Ignoring.")
            transcript = ""

        if not transcript:
            return f"\n\n🎧 AUDIO: {filename}\n[⚠️ Unable to extract meaningful transcript.]"

        # 🧠 Optionally send to Gemini for summarization or enrichment
        prompt_map = {
            "song": "Based on the transcript below, extract complete song lyrics. If possible, also estimate the musical chords or structure (intro, verse, chorus, bridge).",
            "lecture": "Summarize the key points explained in this lecture.",
            "podcast": "List the main discussion topics and insights.",
        }

        log("📊 Classifying audio content...")
        audio_type = classify_video_content_type(transcript, "")  # visuals empty
        log(f">> Detected audio type: {audio_type}")

        user_query = f"What is this audio about?"
        smart_prompt = prompt_map.get(audio_type, f"Based on the transcript, respond to:\n{user_query}")

        log("📨 Sending to Gemini...")
        final_output = summarize_frames_with_gemini("", smart_prompt, transcript)

        # 🎸 Add a helpful tip if Gemini skips the chord part
        if "chord" not in final_output.lower():
            final_output += "\n\n🎸 Tip: For accurate chords, visit Ultimate-Guitar.com or search online for '[song name] chords'."

        # ⛔ Truncate large output to protect frontend
        if len(final_output) > 4000:
            final_output = final_output[:4000] + "...\n[Truncated]"


        log("✅ Transcription and Gemini completed. Returning result...")
        print("📦 Final output length:", len(final_output))

        # return f"\n\n🎧 AUDIO: {filename}\n{final_output}"


        from fastapi.responses import JSONResponse
        return JSONResponse(content={
            "type": "audio",
            "filename": filename,
            "summary": final_output[:4000] + "..." if len(final_output) > 4000 else final_output
        })

        print("✅ Response returned to frontend.")


    except Exception as e:
        logging.error("Audio processing failed: %s", str(e))
        return f"[Error processing audio file {filename}: {e}]"

    finally:
        gc.collect()
        for path in [audio_path, converted_audio_path]:
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                    logging.info(f">> [cleanup] Deleted temp file: {path}")
                except Exception as cleanup_error:
                    logging.warning(f">> [cleanup warning] Failed to delete {path}: {cleanup_error}")

# ============================================# ============================================

def handle_video(filename, decoded_bytes):
    logging.basicConfig(level=logging.INFO)

    def log(msg):
        print(f"[{datetime.now().strftime('%H:%M:%S')}] {msg}")

    video_path = None
    audio_path = None
    converted_audio_path = None

    try:
        log(f"📥 Saving video: {filename}")
        
        # 👇 TEMPORARY DEBUG LINE to test incoming audio
        print("✅ Received decoded_bytes:", len(decoded_bytes), "bytes")
        print("🔍 First few bytes (base64):", b64encode(decoded_bytes[:20]).decode())

        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".mp4") as tmp_video:
            tmp_video.write(decoded_bytes)
            tmp_video.flush()
            video_path = tmp_video.name

        audio_path = video_path.replace(".mp4", ".wav")
        converted_audio_path = video_path.replace(".mp4", "_16k.wav")

        log("🎬 Loading video with MoviePy...")
        clip = VideoFileClip(video_path)

        log("🔊 Extracting audio...")
        clip.audio.write_audiofile(audio_path, logger=None)

        log("✅ Closing MoviePy clip (to avoid WinError 32)...")
        try:
            clip.reader.close()
            if clip.audio:
                clip.audio.reader.close_proc()
        except Exception as cleanup_issue:
            logging.warning(f"Failed to close clip readers: {cleanup_issue}")

        log("🎚 Converting audio to 16kHz mono WAV...")
        # subprocess.run([
        #     "ffmpeg", "-y", "-i", audio_path,
        #     "-ac", "1", "-ar", "16000", "-sample_fmt", "s16",
        #     "-f", "wav", converted_audio_path
        # ], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)

        result = subprocess.run([
            "ffmpeg", "-y", "-i", audio_path,
            "-ac", "1", "-ar", "16000", "-sample_fmt", "s16",
            "-f", "wav", converted_audio_path
        ], capture_output=True, text=True)

        if result.returncode != 0:
            print("❌ FFmpeg failed:", result.stderr)
        else:
            print("✅ FFmpeg success:", result.stdout)


        log("🧠 Loading Whisper model...")
        model = whisper.load_model("base")

        log("🧪 Validating audio...")
        data, samplerate = sf.read(converted_audio_path)
        if len(data) < samplerate:
            raise Exception("Audio too short for transcription")

        log("🗣 Running initial Whisper transcription...")
        result = model.transcribe(converted_audio_path)
        transcript = result.get("text", "").strip()
        lang_code = result.get("language", "en")
        log(f">> Transcript (preview): {transcript[:100]}")

        lang_hint = detect_language_hint(filename, transcript, LANGUAGE_HINTS)
        if lang_hint and lang_hint != lang_code:
            log(f"🗣 Retrying Whisper with language hint: {lang_code} → {lang_hint}")
            try:
                result = model.transcribe(converted_audio_path, language=lang_hint)
                transcript = result.get("text", "").strip()
                lang_code = lang_hint
                log(f">> New transcript (preview): {transcript[:100]}")
            except Exception as e:
                log(f"⚠️ Language override failed: {e}")
                # keep original transcript

        if is_transcript_repetitive(transcript):
            log("⚠️ Transcript is repetitive. Ignoring.")
            transcript = ""

        log("🧠 Performing Gemini visual analysis...")
        visual_summary = analyze_video_frames(video_path, filename, query="Describe what happens visually.")

        if not transcript and not visual_summary:
            return f"🎥 {filename}\n\n⚠️ Unable to extract meaningful content."

        log("📊 Classifying video type...")
        video_type = classify_video_content_type(transcript, visual_summary)
        log(f">> Detected video type: {video_type}")

        prompt_map = {
            "song": "Extract full lyrics from this song video.",
            "cooking": "Summarize the cooking recipe. Mention ingredients and steps.",
            "lecture": "Summarize the key points explained in this lecture.",
            "interview": "List who is speaking and summarize what they say.",
            "vlog": "Describe what the person is doing and where they are.",
        }

        user_query = f"What is this video about?"
        smart_prompt = prompt_map.get(video_type, f"Based on transcript and visuals, respond to:\n{user_query}")
        if not transcript:
            smart_prompt = "Transcript unavailable. Describe what is happening in this video based on visuals only."

        log("📨 Sending to Gemini...")
        final_output = summarize_frames_with_gemini(visual_summary, smart_prompt, transcript)

        log("✅ Transcription and Gemini completed. Returning result...")
        print("📦 Final output length:", len(final_output))

        # return f"\n\n🎥 VIDEO: {filename}\n{final_output}"
        from fastapi.responses import JSONResponse
        return JSONResponse(content={
            "type": "video",
            "filename": filename,
            "summary": final_output[:4000] + "..." if len(final_output) > 4000 else final_output
        })

        print("✅ Response returned to frontend.")


    except Exception as e:
        logging.error("Video processing error: %s", str(e))
        return f"[Error processing video file {filename}: {e}]"

    finally:
        gc.collect()
        for path in [video_path, audio_path, converted_audio_path]:
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                    logging.info(f">> [cleanup] Deleted temp file: {path}")
                except Exception as cleanup_error:
                    logging.warning(f">> [cleanup warning] Failed to delete {path}: {cleanup_error}")

# ===========================================================
# ✅ What You Can Control Easily
# Switch Whisper model: "base" → "medium" or "large"
# Change video resolution: "best[height<=480][ext=mp4]"
# Customize Gemini prompt per content type
# ============================================

LANGUAGE_HINTS = {
    "hi": ["kumar sanu", "bollywood", "naaraaz", "hindi", "sambhala", "mere", "hai", "tum", "dil", "pyaar",
            "asha bhosle", "lata", "arijit", "yeh", "tera", "sapna", "sapne", "zindagi", "mohabbat", "ishq", "shayari",
            "hero", "villain", "gaana", "film", "kahani", "ranbir", "deepika"],
    "ta": ["kollywood", "tamil", "rajini", "vijay", "amma", "enna", "yen", "illa", "thalaiva", "padam",
            "sivakarthikeyan", "vijay sethupathi", "ajith", "kamal", "nayanthara", "thambi", "satham", "kadhal",
            "vettai", "vannakam", "ponniyin", "selvan", "mass", "basha", "veeram"],
    "te": ["tollywood", "telugu", "allu", "mahesh", "raasi", "nuvvu", "vaddu", "padam", "chiranjeevi", "pawan",
            "pushpa", "icon star", "srivalli", "kotha", "bava", "ammo", "veera", "kotha", "nenu", "evaru", "chitti",
            "adavi", "mass", "megastar"],
    "bn": ["bengali", "kolkata", "bangla", "rabindra", "ami", "tumi", "koro", "kotha", "chele", "meyera",
            "song", "gaaner", "sokal", "ratri", "shonar", "bangla", "bijoy", "pran", "anondo", "bhalobasha",
            "bhai", "rong", "misti", "rosogolla"],
    "ml": ["malayalam", "kerala", "mohanlal", "fahadh", "ente", "njan", "alle", "oru", "vannu", "chila",
            "manasil", "amma", "kutty", "mammootty", "nivin", "dileep", "kalyaanam", "pookal", "thaniye", "soorya",
            "thattathin", "marayathe", "kanne"]
}

def detect_language_hint(title: str, transcript: str, lang_hints: dict, threshold=2) -> str:
    text = (title + " " + transcript).lower()
    scores = {lang: sum(1 for word in hints if word in text) for lang, hints in lang_hints.items()}
    best_lang = max(scores, key=scores.get)
    return best_lang if scores[best_lang] >= threshold else ""

def is_transcript_repetitive(text, threshold=5):
    words = text.split()
    unique = set(words)
    return len(unique) < threshold

def handle_youtube_link(url, query):
    video_path = None
    raw_audio_path = None
    converted_audio_path = None

    try:
        print(f"\n>> [yt-dlp] Downloading YouTube video: {url}")
        ydl_opts = {
            "format": "best[height<=480][ext=mp4]",
            "outtmpl": os.path.join(tempfile.gettempdir(), "yt_video.%(ext)s"),
            "merge_output_format": "mp4",
            "quiet": True,
            "verbose": True,
            "noplaylist": True,
        }

        with YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=True)
            video_path = os.path.join(tempfile.gettempdir(), f"yt_video.{info['ext']}")
            title = info.get("title", "Unknown YouTube Video")
        print(f">> [yt-dlp] Video downloaded: {title}")

        print(">> [Whisper] Loading transcription model...")
        model = whisper.load_model("base")

        raw_audio_path = video_path.replace(".mp4", "_raw.wav")
        converted_audio_path = video_path.replace(".mp4", ".wav")

        print(">> [Audio] Extracting and converting to 16kHz mono...")
        with VideoFileClip(video_path) as clip:
            clip.audio.write_audiofile(raw_audio_path, logger=None)

        # subprocess.run([
        #     "ffmpeg", "-y", "-i", raw_audio_path,
        #     "-ac", "1", "-ar", "16000", "-sample_fmt", "s16",
        #     "-f", "wav", converted_audio_path
        # ], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)

        result = subprocess.run([
            "ffmpeg", "-y", "-i", raw_audio_path,
            "-ac", "1", "-ar", "16000", "-sample_fmt", "s16",
            "-f", "wav", converted_audio_path
        ], capture_output=True, text=True)

        if result.returncode != 0:
            print("❌ FFmpeg failed:", result.stderr)
        else:
            print("✅ FFmpeg success:", result.stdout)

        print(">> [Audio] Validating WAV file...")
        data, samplerate = sf.read(converted_audio_path)
        if len(data) < samplerate:
            raise Exception("Audio too short for transcription")

        print(">> [Whisper] First-pass transcription...")
        result = model.transcribe(converted_audio_path)
        transcript = result.get("text", "").strip()
        lang_code = result.get("language", "en")

        print(f">> [Whisper] Detected language: {lang_code}")
        print(f">> [Whisper] Transcript preview: {transcript[:100]}...")

        lang_hint = detect_language_hint(title, transcript, LANGUAGE_HINTS)
        if lang_hint and lang_hint != lang_code:
            print(f">> [Lang Hint] Overriding Whisper language from {lang_code} to {lang_hint}")
            result = model.transcribe(converted_audio_path, language=lang_hint)
            transcript = result.get("text", "").strip()
            lang_code = lang_hint

        if is_transcript_repetitive(transcript):
            print(">> [Whisper] Transcript appears repetitive or invalid. Using visual-only fallback.")
            transcript = ""

        print(">> [Gemini] Performing visual analysis...")
        visual_summary = analyze_video_frames(video_path, title, query="What do you see happening visually in this video?")

        if not transcript and not visual_summary:
            return f"🎥 {title}\n\n⚠️ Unable to extract meaningful content from this video."

        print(">> [Gemini] Classifying content type...")
        video_type = classify_video_content_type(transcript, visual_summary)
        print(f">> [Gemini] Detected video type: {video_type}")

        prompt_map = {
            "song": "Extract full lyrics from this song video.",
            "cooking": "Summarize the cooking recipe. Mention ingredients and steps.",
            "lecture": "Summarize the key points explained in this lecture.",
            "interview": "List who is speaking and summarize what they say.",
            "vlog": "Describe what the person is doing and where they are."
        }
        smart_prompt = prompt_map.get(video_type, f"Based on the transcript and visuals, respond to:\n{query}")
        if not transcript:
            smart_prompt = "The transcript was unreliable. Based on visuals only, describe what is happening in this video."

        print(f">> [Gemini] Final prompt:\n{smart_prompt}")
        final_output = summarize_frames_with_gemini(visual_summary, smart_prompt, transcript)

        return f"🎥 YouTube: {title}\n\n{final_output}"

    except Exception as e:
        print(">> [Error Traceback]")
        print(traceback.format_exc())
        return f"[ERROR during YouTube video processing: {e}]"

    finally:
        gc.collect()
        for path in [video_path, raw_audio_path, converted_audio_path]:
            if path and os.path.exists(path):
                try:
                    os.remove(path)
                    print(f">> [cleanup] Deleted temp file: {path}")
                except Exception as cleanup_error:
                    print(f">> [cleanup warning] Failed to delete {path}: {cleanup_error}")

# ===========================================================
def handle_pdf(filename, decoded_bytes):
    try:
        pdf_reader = PdfReader(io.BytesIO(decoded_bytes))
        text = "\n".join(page.extract_text() for page in pdf_reader.pages if page.extract_text())
        return f"\n\n📄 PDF: {filename}\n{text.strip()[:2000]}"
    except Exception as e:
        return f"[Error reading {filename}: {e}]"

# ============================================
def handle_csv(filename, decoded_bytes):
    try:
        csv_text = io.StringIO(decoded_bytes.decode("utf-8"))
        reader = csv.reader(csv_text)
        preview_rows = [", ".join(row) for _, row in zip(range(10), reader)]
        return f"\n\n📊 CSV: {filename}\n" + "\n".join(preview_rows)
    except Exception as e:
        return f"[Error reading {filename}: {e}]"
# ============================================
def handle_docx(filename, decoded_bytes):
    try:
        doc = Document(io.BytesIO(decoded_bytes))
        text = "\n".join([p.text for p in doc.paragraphs])
        return f"\n\n📝 DOCX: {filename}\n{text.strip()[:2000]}"
    except Exception as e:
        return f"[Error reading {filename}: {e}]"
# ============================================
def handle_pptx(filename, decoded_bytes):
    try:
        ppt = Presentation(io.BytesIO(decoded_bytes))
        text = "\n".join([shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")])
        return f"\n\n📽 PPTX: {filename}\n{text.strip()[:2000]}"
    except Exception as e:
        return f"[Error reading {filename}: {e}]"
# ============================================
def handle_xlsx(filename, decoded_bytes):
    try:
        wb = load_workbook(io.BytesIO(decoded_bytes), read_only=True)
        sheet = wb.active
        rows = [[str(cell.value) for cell in row] for row in sheet.iter_rows(min_row=1, max_row=10)]
        preview = "\n".join([", ".join(row) for row in rows])
        return f"\n\n📊 XLSX: {filename}\n{preview}"
    except Exception as e:
        return f"[Error reading {filename}: {e}]"
# ============================================
def handle_json(filename, decoded_bytes):
    try:
        parsed = json.loads(decoded_bytes.decode("utf-8"))
        return f"\n\n🧾 JSON: {filename}\n{json.dumps(parsed, indent=2)[:2000]}"
    except Exception as e:
        return f"[Error reading {filename}: {e}]"
# ============================================
def handle_xml(filename, decoded_bytes):
    try:
        root = ET.fromstring(decoded_bytes.decode("utf-8"))
        return f"\n\n🧾 XML: {filename}\n{ET.tostring(root, encoding='unicode')[:2000]}"
    except Exception as e:
        return f"[Error reading {filename}: {e}]"
# ============================================
def handle_archive(filename, decoded_bytes):
    try:
        content_list = []
        if zipfile.is_zipfile(io.BytesIO(decoded_bytes)):
            with zipfile.ZipFile(io.BytesIO(decoded_bytes)) as zf:
                content_list = zf.namelist()
        elif tarfile.is_tarfile(io.BytesIO(decoded_bytes)):
            with tarfile.open(fileobj=io.BytesIO(decoded_bytes)) as tf:
                content_list = tf.getnames()
        return f"\n\n🗂 Archive: {filename}\nContains files:\n" + "\n".join(content_list[:10])
    except Exception as e:
        return f"[Error reading archive {filename}: {e}]"
# ===========================================================
# END
# ===========================================================


